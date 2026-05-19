from __future__ import annotations

import argparse
import re
from collections import defaultdict
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable, Sequence

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_INPUT_XLSX = Path(r"C:\Users\rober\Downloads\oz_dashboard_goods_merged.xlsx")
DEFAULT_TEMPLATE_PPTX = Path(r"C:\Users\rober\Downloads\01.04.2026 (1).pptx")
DEFAULT_OUTPUT_PPTX = Path(r"C:\Users\rober\Downloads\oz_dashboard_goods_generated_2026-04-14_stocks-fixed.pptx")

TARGET_SKUS = [
    "LL-H4-BBox",
    "LL-H7-BBox",
    "LL-H11-BBox",
    "LL-H1-BBox",
    "LL-HB3-BBox",
    "BIL-M1",
    "x-pro-2.5",
    "g6-x-l-2.5",
    "GLAZKI-2.5",
    "GLAZKI-3.0",
    "2XL-D1S-5000K",
    "2XL-D2S-5000K",
    "2XL-D3S-5000K",
    "2XL-D4S-5000K",
    "2XL-D1S-6000K",
    "2XL-D2S-6000K",
    "2XL-D3S-6000K",
    "2XL-D4S-6000K",
    "TU-X-VST-1",
    "TU-X-VST-2",
    "TU-X-TOY-1",
    "TU-X-TOY-2",
    "TU-X-GRT-1",
    "TU-X-GRT-2",
    "TU-X-PRA-2",
    "TU-X-PRA-1",
    "TU-X-NS-1",
    "TU-X-NS-2",
    "A13-30-5500",
    "A3MAX-2.5",
    "A3MAX-3.0",
    "AOZDXO1156",
    "Dragon",
    "PROVOD-L-H4",
    "BIL-D5",
    "K7",
    "MINILENS-F50-H4",
    "BIL-D10",
    "BIL-D6",
    "LL-H7-BULL",
    "LL-H11-BULL",
    "LL-H4-BULL",
    "LL-H1-BULL",
    "LL-HB3-BULL",
    "GERM-1",
    "MASK-L-101",
    "MASK-L-102",
    "MASK-L-103",
    "MASK-L-104",
    "MASK-L-105",
    "LL-HYBRID-D1S",
    "LL-HYBRID-D3S",
    "LL-H7-RR",
    "LL-H4-Y7D",
    "LL-H4-LIONS",
    "LL-HB3-LIONS",
    "LL-HB4-LIONS",
    "LL-H1-LIONS",
    "LL-H7-LIONS",
    "LL-H11-LIONS",
    "LL-HALOGEN-H4",
    "LL-HALOGEN-H7",
    "LL-HALOGEN-H11",
    "LL-COMPACT-H1",
    "LL-COMPACT-H11",
    "LL-COMPACT-H4",
    "LL-COMPACT-H7",
]

MANUAL_STOCK_LINES = [
    "LL-COMPACT-H7-\t13",
    "LL-COMPACT-H11-\t21",
    "LL-H4-BBox\t12283",
    "x-pro-2.5\t1609",
    "LL-COMPACT-H1-\t95",
    "LL-COMPACT-H4-\t34",
    "BIL-M1\t4883",
    "LL-H7-LIONS\t1516",
    "LL-HYBRID-D1S\t748",
    "LL-H4-Y7D\t662",
    "LL-H7-BBox\t8090",
    "Dragon\t70",
    "TU-X-VST-2\t358",
    "LL-HYBRID-D3S\t756",
    "TU-X-VST-1\t1626",
    "LL-H11-LIONS\t811",
    "LL-H7-BULL\t3801",
    "LL-H4-LIONS\t973",
    "LL-H1-LIONS\t422",
    "K7\t102",
    "LL-HALOGEN-H7\t274",
    "LL-H4-BULL\t2590",
    "TU-X-PRA-2\t62",
    "A13-30-5500\t178",
    "TU-X-TOY-1\t420",
    "MASK-L-105\t651",
    "LL-HB3-LIONS\t214",
    "LL-H1-BBox\t4305",
    "TU-X-TOY-2\t0",
    "MASK-L-101\t1306",
    "TU-X-GRT-1\t324",
    "2XL-D1S-5000K\t74",
    "GERM-1\t866",
    "PROVOD-L-H4\t3307",
    "MASK-L-103\t476",
    "LL-H11-BBox\t1049",
    "LL-HB3-BBox\t622",
    "LL-H1-BULL\t177",
    "2XL-D2S-6000K\t213",
    "2XL-D4S-6000K\t233",
    "2XL-D3S-5000K\t44",
    "LL-HALOGEN-H11\t354",
    "AOZDXO1156\t146",
    "TU-X-PRA-1\t1",
    "MASK-L-104\t716",
    "LL-HALOGEN-H4\t0",
    "2XL-D1S-6000K\t1117",
    "LL-HB4-LIONS\t183",
    "TU-X-GRT-2\t410",
    "LL-H11-BULL\t2",
    "2XL-D3S-6000K\t347",
    "MASK-L-102\t324",
    "2XL-D4S-5000K\t131",
    "LL-H7-RR\t811",
    "g6-x-l-2.5\t2145",
    "GLAZKI-3.0\t598",
    "2XL-D2S-5000K\t265",
    "GLAZKI-2.5\t1008",
    "A3MAX-2.5\t7",
    "2XL-D2S-6000K\t213",
    "LL-HALOGEN-H7\t274",
    "LL-COMPACT-H4-\t34",
    "LL-HALOGEN-H4\t0",
    "LL-H4-BBox\t12283",
    "TU-X-NS-1\t0",
    "MINILENS-F50-H4\t0",
]

METRIC_DISPLAY_NAMES = [
    "Заказы",
    "Продажи",
    "DRR %",
    "ROCE %",
    "Чистая прибыль",
    "Прибыль/шт",
    "Ср. цена",
]

BLUE = RGBColor(0x2C, 0x5E, 0x92)
ZEBRA = RGBColor(0xED, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)


@dataclass(frozen=True)
class PeriodInfo:
    title: str
    start: datetime
    end: datetime
    days: int


@dataclass
class AggregatedMetrics:
    orders: float = 0.0
    sales_units: float = 0.0
    revenue: float = 0.0
    profit: float = 0.0
    stock: float = 0.0
    drr_weighted_sum: float = 0.0
    roce_weighted_sum: float = 0.0
    weight_sum: float = 0.0
    row_count: int = 0

    def add_row(
        self,
        *,
        orders: float,
        sales_units: float,
        revenue: float,
        profit: float,
        stock: float,
        drr: float | None,
        roce: float | None,
    ) -> None:
        self.orders += orders
        self.sales_units += sales_units
        self.revenue += revenue
        self.profit += profit
        self.stock += stock
        weight = revenue if revenue > 0 else (sales_units if sales_units > 0 else 1.0)
        if drr is not None:
            self.drr_weighted_sum += drr * weight
        if roce is not None:
            self.roce_weighted_sum += roce * weight
        self.weight_sum += weight
        self.row_count += 1

    @property
    def drr(self) -> float | None:
        if self.weight_sum == 0:
            return None
        return self.drr_weighted_sum / self.weight_sum

    @property
    def roce(self) -> float | None:
        if self.weight_sum == 0:
            return None
        return self.roce_weighted_sum / self.weight_sum

    @property
    def profit_per_unit(self) -> float | None:
        if self.sales_units == 0:
            return None
        return self.profit / self.sales_units

    @property
    def avg_price(self) -> float | None:
        if self.sales_units == 0:
            return None
        return self.revenue / self.sales_units

    @property
    def has_data(self) -> bool:
        return self.row_count > 0


@dataclass
class StockInfo:
    available: float = 0.0
    incoming: float = 0.0


def normalize_header(value: str) -> str:
    return re.sub(r"\s+", " ", value.strip().casefold())


def normalize_sku(value: str) -> str:
    text = str(value or "").strip()
    text = re.sub(r"[-.\s]+$", "", text)
    text = re.sub(r"^led(?=-)", "ll", text, flags=re.IGNORECASE)
    return text.casefold()


def safe_number(value) -> float:
    if value in (None, ""):
        return 0.0
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip().replace("\xa0", " ").replace(" ", "")
    text = text.replace(",", ".")
    if not text:
        return 0.0
    try:
        return float(text)
    except ValueError:
        return 0.0


def parse_period_title(title: str) -> PeriodInfo | None:
    match = re.fullmatch(r"\s*(\d{2}\.\d{2}\.\d{4})\s*-\s*(\d{2}\.\d{2}\.\d{4})\s*", title)
    if not match:
        return None
    start = datetime.strptime(match.group(1), "%d.%m.%Y")
    end = datetime.strptime(match.group(2), "%d.%m.%Y")
    return PeriodInfo(title=title, start=start, end=end, days=(end - start).days + 1)


def find_first(headers: list[str], keyword_groups: Iterable[tuple[str, ...]], *, exclude: Iterable[int] = ()) -> int:
    excluded = set(exclude)
    normalized = [normalize_header(h) for h in headers]
    for keywords in keyword_groups:
        for idx, header in enumerate(normalized):
            if idx in excluded:
                continue
            if all(keyword in header for keyword in keywords):
                return idx
    raise KeyError(f"Column not found for keywords: {keyword_groups}")


def build_column_map(headers: list[str]) -> dict[str, int]:
    orders_idx = find_first(headers, [("заказы, шт",), ("заказы",)])
    sales_idx = find_first(
        headers,
        [("количество продаж",), ("продажи, шт",), ("units",), ("шт",)],
        exclude=[orders_idx],
    )
    revenue_idx = find_first(headers, [("сумма продаж",), ("выручка",)])
    return {
        "sku_store": find_first(headers, [("sku магазина",), ("артикул",), ("sku",)]),
        "sku_ozon": find_first(headers, [("sku ozon",), ("sku",)], exclude=[]),
        "orders": orders_idx,
        "sales_units": sales_idx,
        "revenue": revenue_idx,
        "drr": find_first(headers, [("tacoo",), ("дрд",), ("дрр",)]),
        "roce": find_first(headers, [("roce",)]),
        "profit": find_first(headers, [("чистая прибыль",), ("profit",)]),
        "stock": find_first(headers, [("доступно к продаже",), ("остаток",)]),
    }


def build_column_map(headers: list[str]) -> dict[str, int]:
    orders_idx = find_first(headers, [("Р·Р°РєР°Р·С‹, С€С‚",), ("Р·Р°РєР°Р·С‹",)])
    sales_idx = find_first(
        headers,
        [("РєРѕР»РёС‡РµСЃС‚РІРѕ РїСЂРѕРґР°Р¶",), ("РїСЂРѕРґР°Р¶Рё, С€С‚",), ("units",), ("С€С‚",)],
        exclude=[orders_idx],
    )
    revenue_idx = find_first(headers, [("СЃСѓРјРјР° РїСЂРѕРґР°Р¶",), ("РІС‹СЂСѓС‡РєР°",)])
    sku_store_idx = find_first(headers, [("sku РјР°РіР°Р·РёРЅР°",), ("Р°СЂС‚РёРєСѓР»",), ("sku",)])

    normalized = [normalize_header(h) for h in headers]
    sku_ozon_idx = None
    for idx, header in enumerate(normalized):
        if idx == sku_store_idx:
            continue
        if "sku ozon" in header or "sku wb" in header or "id карточки" in header:
            sku_ozon_idx = idx
            break
    if sku_ozon_idx is None:
        for idx, header in enumerate(normalized):
            if idx != sku_store_idx and "sku" in header:
                sku_ozon_idx = idx
                break
    if sku_ozon_idx is None:
        sku_ozon_idx = sku_store_idx

    return {
        "sku_store": sku_store_idx,
        "sku_ozon": sku_ozon_idx,
        "orders": orders_idx,
        "sales_units": sales_idx,
        "revenue": revenue_idx,
        "drr": find_first(headers, [("tacoo",), ("РґСЂРґ",), ("РґСЂСЂ",)]),
        "roce": find_first(headers, [("roce",)]),
        "profit": find_first(headers, [("С‡РёСЃС‚Р°СЏ РїСЂРёР±С‹Р»СЊ",), ("profit",)]),
        "stock": find_first(headers, [("РґРѕСЃС‚СѓРїРЅРѕ Рє РїСЂРѕРґР°Р¶Рµ",), ("РѕСЃС‚Р°С‚РѕРє",)]),
    }


def build_column_map(headers: list[str]) -> dict[str, int]:
    normalized = [normalize_header(h) for h in headers]

    def find_idx(candidates: list[str], *, exclude: set[int] | None = None) -> int:
        excluded = exclude or set()
        for candidate in candidates:
            for idx, header in enumerate(normalized):
                if idx in excluded:
                    continue
                if candidate in header:
                    return idx
        raise KeyError(f"Column not found for candidates: {candidates}")

    sku_store_idx = find_idx(
        [
            "sku \u043c\u0430\u0433\u0430\u0437\u0438\u043d\u0430",
            "\u0430\u0440\u0442\u0438\u043a\u0443\u043b \u043c\u0430\u0433\u0430\u0437\u0438\u043d\u0430",
            "\u0430\u0440\u0442\u0438\u043a\u0443\u043b",
        ]
    )
    sku_market_idx = None
    for candidate in ["sku ozon", "sku wb", "id \u043a\u0430\u0440\u0442\u043e\u0447\u043a\u0438", "sku"]:
        for idx, header in enumerate(normalized):
            if idx == sku_store_idx:
                continue
            if candidate in header:
                sku_market_idx = idx
                break
        if sku_market_idx is not None:
            break
    if sku_market_idx is None:
        sku_market_idx = sku_store_idx

    orders_idx = find_idx(["\u0437\u0430\u043a\u0430\u0437\u044b"])
    sales_idx = find_idx(
        [
            "\u043a\u043e\u043b-\u0432\u043e \u043f\u0440\u043e\u0434\u0430\u0436",
            "\u043f\u0440\u043e\u0434\u0430\u0436\u0438, \u0448\u0442",
            "units",
        ],
        exclude={orders_idx},
    )
    revenue_idx = find_idx(
        [
            "\u0441\u0443\u043c\u043c\u0430 \u043f\u0440\u043e\u0434\u0430\u0436",
            "\u0432\u044b\u0440\u0443\u0447\u043a\u0430",
        ]
    )

    return {
        "sku_store": sku_store_idx,
        "sku_ozon": sku_market_idx,
        "orders": orders_idx,
        "sales_units": sales_idx,
        "revenue": revenue_idx,
        "drr": find_idx(["tacoo", "\u0434\u0440\u0440", "\u0434\u0440\u0434"]),
        "roce": find_idx(["roce"]),
        "profit": find_idx(["\u0447\u0438\u0441\u0442\u0430\u044f \u043f\u0440\u0438\u0431\u044b\u043b\u044c", "profit"]),
        "stock": find_idx(
            [
                "\u043f\u043e\u043b\u043d\u044b\u0439 \u043e\u0441\u0442\u0430\u0442\u043e\u043a",
                "\u0434\u043e\u0441\u0442\u0443\u043f\u043d\u043e \u043a \u043f\u0440\u043e\u0434\u0430\u0436\u0435",
                "\u043e\u0441\u0442\u0430\u0442\u043e\u043a",
            ]
        ),
    }


def choose_weekly_periods(periods: list[PeriodInfo]) -> list[PeriodInfo]:
    weekly = [p for p in periods if 6 <= p.days <= 8]
    if len(weekly) >= 4:
        return sorted(weekly, key=lambda p: (p.start, p.end))[:4]
    remaining = [p for p in periods if p not in weekly]
    remaining.sort(key=lambda p: (abs(p.days - 7), p.start, p.end))
    needed = 4 - len(weekly)
    weekly.extend(remaining[:needed])
    return sorted(weekly, key=lambda p: (p.start, p.end))


def format_int_like(value: float | None) -> str:
    if value is None:
        return "—"
    return f"{int(round(value)):,}".replace(",", " ")


def format_days_coverage(stock_value: float, sales_28d: float) -> str:
    if sales_28d <= 0:
        return "—"
    return str(int(round(stock_value * 28 / sales_28d)))


def format_percent(value: float | None) -> str:
    if value is None:
        return "—"
    return f"{int(round(value))}%"


def format_plain_number(value: float | None) -> str:
    if value is None:
        return "—"
    rounded = int(round(value))
    sign = "+" if rounded > 0 else ""
    return f"{sign}{rounded:,}".replace(",", " ")


def format_pp_change(value: float | None) -> str:
    if value is None:
        return "—"
    rounded = int(round(value))
    sign = "+" if rounded > 0 else ""
    return f"{sign}{rounded} п.п."


def format_percent_change(old: float | None, new: float | None) -> str:
    if old in (None, 0) or new is None:
        return "—"
    change = (new - old) / old * 100
    return f"{change:+.1f}%"


def metric_value(metrics: AggregatedMetrics, metric: str) -> float | None:
    if metric == "Заказы":
        return metrics.orders
    if metric == "Продажи":
        return metrics.sales_units
    if metric == "DRR %":
        return metrics.drr
    if metric == "ROCE %":
        return metrics.roce
    if metric == "Чистая прибыль":
        return metrics.profit
    if metric == "Прибыль/шт":
        return metrics.profit_per_unit
    if metric == "Ср. цена":
        return metrics.avg_price
    raise KeyError(metric)


def format_metric(metric: str, value: float | None) -> str:
    if metric in {"Заказы", "Продажи", "Чистая прибыль", "Прибыль/шт", "Ср. цена"}:
        return format_int_like(value)
    if metric in {"DRR %", "ROCE %"}:
        return format_percent(value)
    raise KeyError(metric)


def format_change(metric: str, old: float | None, new: float | None) -> str:
    if old is None or new is None:
        return "—"
    diff = new - old
    if metric in {"DRR %", "ROCE %"}:
        return format_pp_change(diff)
    return format_plain_number(diff)


def get_template_size(template_pptx: Path) -> tuple[int, int]:
    prs = Presentation(str(template_pptx))
    return prs.slide_width, prs.slide_height


def parse_args(argv: Sequence[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate an Ozon SKU PowerPoint presentation from a merged Excel workbook.",
    )
    parser.add_argument(
        "--input-xlsx",
        type=Path,
        default=DEFAULT_INPUT_XLSX,
        help=f"Source workbook path. Default: {DEFAULT_INPUT_XLSX}",
    )
    parser.add_argument(
        "--template-pptx",
        type=Path,
        default=DEFAULT_TEMPLATE_PPTX,
        help=f"Template presentation path used only for slide size. Default: {DEFAULT_TEMPLATE_PPTX}",
    )
    parser.add_argument(
        "--output-pptx",
        type=Path,
        default=DEFAULT_OUTPUT_PPTX,
        help=f"Generated presentation path. Default: {DEFAULT_OUTPUT_PPTX}",
    )
    parser.add_argument(
        "--stocks-xlsx",
        type=Path,
        default=None,
        help="Optional stocks workbook with SKU, available stock, and incoming stock columns.",
    )
    return parser.parse_args(argv)


def build_manual_stock_overrides() -> dict[str, StockInfo]:
    stock_by_sku: dict[str, StockInfo] = {}
    for line in MANUAL_STOCK_LINES:
        sku_raw, qty_raw = line.split("\t", 1)
        stock_by_sku[normalize_sku(sku_raw)] = StockInfo(available=safe_number(qty_raw), incoming=0.0)
    return dict(stock_by_sku)


def load_stock_overrides(stocks_xlsx: Path | None) -> dict[str, StockInfo]:
    if stocks_xlsx is None:
        return build_manual_stock_overrides()

    workbook = load_workbook(stocks_xlsx, read_only=True, data_only=True)
    ws = workbook[workbook.sheetnames[0]]
    rows = ws.iter_rows(values_only=True)
    headers = [str(value).strip() if value is not None else "" for value in next(rows)]
    column_map = {
        "sku_store": find_first(headers, [("sku продавца",), ("sku магазина",), ("артикул",), ("sku",)]),
        "available": find_first(headers, [("доступно к продаже",), ("доступных к продаже",)]),
        "incoming": find_first(headers, [("ожидается поставка",)]),
    }

    stock_by_sku: dict[str, StockInfo] = {}
    for row in rows:
        sku = normalize_sku(row[column_map["sku_store"]])
        if not sku:
            continue
        current = stock_by_sku.setdefault(sku, StockInfo())
        current.available += safe_number(row[column_map["available"]])
        current.incoming += safe_number(row[column_map["incoming"]])
    return stock_by_sku


def set_text_style(cell, *, align: PP_ALIGN, color: RGBColor, bold: bool = False, size_pt: int = 9) -> None:
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color


def style_table(shape, left_align_cols: set[int]) -> None:
    table = shape.table
    rows = len(table.rows)
    cols = len(table.columns)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                align = PP_ALIGN.LEFT if c in left_align_cols else PP_ALIGN.RIGHT
                set_text_style(cell, align=align, color=WHITE, bold=True)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ZEBRA if r % 2 == 1 else WHITE
                align = PP_ALIGN.LEFT if c in left_align_cols else PP_ALIGN.RIGHT
                set_text_style(cell, align=align, color=BLACK, bold=False)


def fill_table(shape, data: list[list[str]], left_align_cols: set[int], col_widths_inches: list[float]) -> None:
    table = shape.table
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    for idx, width in enumerate(col_widths_inches):
        table.columns[idx].width = Inches(width)
    style_table(shape, left_align_cols)


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(9.1), Inches(0.75))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = BLACK


def create_no_data_slide(prs: Presentation, sku: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, f"{sku} — нет данных")


def create_data_slide(
    prs: Presentation,
    sku: str,
    weekly_periods: list[PeriodInfo],
    big_periods: list[PeriodInfo],
    metrics_by_period: dict[str, dict[str, AggregatedMetrics]],
    last_period: PeriodInfo,
    stock_overrides: dict[str, StockInfo],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    latest_metrics = metrics_by_period[last_period.title].get(sku)
    new_big_metrics = metrics_by_period[big_periods[1].title].get(sku)
    normalized_sku = normalize_sku(sku)
    if normalized_sku in stock_overrides:
        stock_info = stock_overrides[normalized_sku]
    else:
        stock_info = StockInfo()
    if normalized_sku not in stock_overrides and latest_metrics and latest_metrics.has_data:
        stock_info = StockInfo(available=latest_metrics.stock, incoming=0.0)
    sales_new = new_big_metrics.sales_units if new_big_metrics and new_big_metrics.has_data else 0.0
    available_days = format_days_coverage(stock_info.available, sales_new)
    incoming_days = format_days_coverage(stock_info.incoming, sales_new)
    add_title(
        slide,
        f"{sku} | Остаток: {format_int_like(stock_info.available)} ({available_days} д.) + поставка: {format_int_like(stock_info.incoming)} ({incoming_days} д.)",
    )

    weekly_rows = [["Период", *METRIC_DISPLAY_NAMES]]
    for period in weekly_periods:
        metrics = metrics_by_period[period.title].get(sku, AggregatedMetrics())
        weekly_rows.append(
            [
                period.title,
                format_metric("Заказы", metric_value(metrics, "Заказы")),
                format_metric("Продажи", metric_value(metrics, "Продажи")),
                format_metric("DRR %", metric_value(metrics, "DRR %")),
                format_metric("ROCE %", metric_value(metrics, "ROCE %")),
                format_metric("Чистая прибыль", metric_value(metrics, "Чистая прибыль")),
                format_metric("Прибыль/шт", metric_value(metrics, "Прибыль/шт")),
                format_metric("Ср. цена", metric_value(metrics, "Ср. цена")),
            ]
        )

    weekly_shape = slide.shapes.add_table(5, 8, Inches(0.3), Inches(1.35), Inches(9.4), Inches(2.2))
    fill_table(
        weekly_shape,
        weekly_rows,
        left_align_cols={0},
        col_widths_inches=[2.45, 0.8, 0.8, 0.7, 0.7, 1.3, 0.85, 0.8],
    )

    old_metrics = metrics_by_period[big_periods[0].title].get(sku, AggregatedMetrics())
    new_metrics = metrics_by_period[big_periods[1].title].get(sku, AggregatedMetrics())
    compare_rows = [["Метрика", big_periods[0].title, big_periods[1].title, "Изм.", "Изм., %"]]
    for metric in METRIC_DISPLAY_NAMES:
        old_value = metric_value(old_metrics, metric)
        new_value = metric_value(new_metrics, metric)
        compare_rows.append(
            [
                metric,
                format_metric(metric, old_value),
                format_metric(metric, new_value),
                format_change(metric, old_value, new_value),
                format_percent_change(old_value, new_value),
            ]
        )

    compare_shape = slide.shapes.add_table(8, 5, Inches(0.3), Inches(3.75), Inches(9.4), Inches(3.1))
    fill_table(
        compare_shape,
        compare_rows,
        left_align_cols={0},
        col_widths_inches=[1.75, 2.2, 2.2, 1.0, 1.1],
    )


def main(argv: Sequence[str] | None = None) -> None:
    args = parse_args(argv)
    input_xlsx = args.input_xlsx.expanduser()
    template_pptx = args.template_pptx.expanduser()
    output_pptx = args.output_pptx.expanduser()
    stocks_xlsx = args.stocks_xlsx.expanduser() if args.stocks_xlsx else None

    workbook = load_workbook(input_xlsx, read_only=True, data_only=True)
    period_sheets: list[tuple[PeriodInfo, str]] = []
    for sheet_name in workbook.sheetnames:
        parsed = parse_period_title(sheet_name)
        if parsed:
            period_sheets.append((parsed, sheet_name))

    if len(period_sheets) < 2:
        raise RuntimeError("Not enough period sheets found in workbook.")

    period_sheets.sort(key=lambda item: (item[0].start, item[0].end))
    periods = [item[0] for item in period_sheets]
    big_periods = sorted(sorted(periods, key=lambda p: (-p.days, p.start))[:2], key=lambda p: p.start)
    weekly_periods = choose_weekly_periods(periods)
    last_period = max(periods, key=lambda p: p.end)

    sku_lookup = {normalize_sku(sku): sku for sku in TARGET_SKUS}
    metrics_by_period: dict[str, dict[str, AggregatedMetrics]] = {
        period.title: defaultdict(AggregatedMetrics) for period in periods
    }
    stock_overrides = load_stock_overrides(stocks_xlsx)

    for period, sheet_name in period_sheets:
        ws = workbook[sheet_name]
        rows = ws.iter_rows(values_only=True)
        headers = [str(value).strip() if value is not None else "" for value in next(rows)]
        column_map = build_column_map(headers)
        for row in rows:
            raw_store_sku = row[column_map["sku_store"]]
            raw_ozon_sku = row[column_map["sku_ozon"]]
            matched = None
            for candidate in (raw_store_sku, raw_ozon_sku):
                normalized = normalize_sku(candidate)
                if normalized in sku_lookup:
                    matched = sku_lookup[normalized]
                    break
            if not matched:
                continue
            metrics_by_period[period.title][matched].add_row(
                orders=safe_number(row[column_map["orders"]]),
                sales_units=safe_number(row[column_map["sales_units"]]),
                revenue=safe_number(row[column_map["revenue"]]),
                profit=safe_number(row[column_map["profit"]]),
                stock=safe_number(row[column_map["stock"]]),
                drr=safe_number(row[column_map["drr"]]) if row[column_map["drr"]] not in (None, "") else None,
                roce=safe_number(row[column_map["roce"]]) if row[column_map["roce"]] not in (None, "") else None,
            )

    prs = Presentation()
    slide_width, slide_height = get_template_size(template_pptx)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    for sku in TARGET_SKUS:
        has_any = any(metrics_by_period[period.title].get(sku, AggregatedMetrics()).has_data for period in periods)
        if not has_any:
            create_no_data_slide(prs, sku)
            continue
        create_data_slide(prs, sku, weekly_periods, big_periods, metrics_by_period, last_period, stock_overrides)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)

    print(f"Saved: {output_pptx}")
    print(f"Big periods: {[p.title for p in big_periods]}")
    print(f"Weekly periods: {[p.title for p in weekly_periods]}")
    print(f"Last period: {last_period.title}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
